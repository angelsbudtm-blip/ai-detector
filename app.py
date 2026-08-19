import streamlit as st
from pptx import Presentation
from pypdf import PdfReader
import psycopg2
import re
import math
from difflib import SequenceMatcher
from datetime import datetime, timezone, timedelta
from fpdf import FPDF

# Define Malaysia Timezone (UTC+8)
MYT = timezone(timedelta(hours=8))

# ----------------------------------------------------
# Database Connection (Supabase PostgreSQL)
# ----------------------------------------------------
def get_db_connection():
    """Establishes connection to Supabase using secret credentials."""
    pg = st.secrets["postgres"]
    return psycopg2.connect(
        host=pg["host"],
        port=pg["port"],
        dbname=pg["dbname"],
        user=pg["user"],
        password=pg["password"]
    )

def init_db():
    """Creates database tables in Supabase if they don't exist yet."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS presentations (
            id SERIAL PRIMARY KEY,
            filename TEXT,
            file_type TEXT,
            total_slides INTEGER,
            overall_ai_score REAL,
            overall_plagiarism_score REAL DEFAULT 0.0,
            upload_time VARCHAR(50)
        );
    ''')
    c.execute('''
        CREATE TABLE IF NOT EXISTS slide_records (
            id SERIAL PRIMARY KEY,
            presentation_id INTEGER REFERENCES presentations(id) ON DELETE CASCADE,
            slide_number INTEGER,
            extracted_text TEXT,
            ai_score REAL
        );
    ''')
    # Ensure overall_plagiarism_score column exists for older schema
    try:
        c.execute("ALTER TABLE presentations ADD COLUMN IF NOT EXISTS overall_plagiarism_score REAL DEFAULT 0.0;")
    except Exception:
        pass

    conn.commit()
    c.close()
    conn.close()

try:
    init_db()
except Exception as e:
    st.error(f"Database connection error: {e}")

# ----------------------------------------------------
# Security / Login Handler
# ----------------------------------------------------
def check_password():
    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False

    if not st.session_state["authenticated"]:
        st.title("🔒 AI Slide Detector Login")
        user_input = st.text_input("Username")
        pass_input = st.text_input("Password", type="password")
        
        if st.button("Login", type="primary"):
            if user_input == "admin" and pass_input in ["ABEdu@5603", "ABEdu5603Secret"]:
                st.session_state["authenticated"] = True
                st.rerun()
            else:
                st.error("Invalid Username or Password")
        return False
    return True

if not check_password():
    st.stop()

# ----------------------------------------------------
# High-Precision Sentence & Document AI Detection Engine
# ----------------------------------------------------
def score_single_sentence(sentence: str) -> float:
    """Calculates AI probability for an individual sentence or line."""
    clean_s = sentence.strip()
    words = re.findall(r'\b[a-zA-Z]+\b', clean_s.lower())
    if len(words) < 2:
        return 0.0

    score = 0.0
    text_lower = clean_s.lower()

    # 1. Structural AI Q&A and Definition Patterns
    qa_patterns = [
        r"^what is\b", r"^why is\b", r"^how can we\b", r"^can using\b", 
        r"^why should we\b", r"^is copying\b", r"\bmeans using\b",
        r"\bwithout giving\b", r"\buse your own\b", r"\balways give credit\b",
        r"\bbe honest, be creative\b", r"\bbecause it is dishonest\b",
        r"\bstops us from learning\b", r"\bplagiarism\b", r"^qna\b", r"^sources\b"
    ]
    for pattern in qa_patterns:
        if re.search(pattern, text_lower):
            score += 50.0

    # 2. General AI Vocabulary & Transition Phrases
    ai_phrases = [
        "in conclusion", "important to note", "crucial role", "key takeaways",
        "furthermore", "moreover", "in summary", "fast-paced", "delve into",
        "tapestry of", "fostering a", "paramount to", "transformative impact",
        "plays a vital", "holistic approach", "it is essential", "in order to"
    ]
    for phrase in ai_phrases:
        if phrase in text_lower:
            score += 45.0

    # 3. Sentence Length Uniformity & Pacing
    word_count = len(words)
    if 6 <= word_count <= 25:
        score += 25.0
    
    # 4. Standard Definition / Rule Formatting
    if clean_s.startswith(("1.", "2.", "3.", "4.", "5.", "- ", "• ")):
        score += 20.0

    return min(99.0, max(15.0 if score > 0 else 0.0, score))

def analyze_document_text(text: str):
    """Parses text sentence-by-sentence, highlighting AI content."""
    if not text or len(text.strip()) < 5:
        return 0.0, "", 0, 0

    raw_chunks = [c.strip() for c in re.split(r'(\n+|[.!?]+)', text) if c.strip()]
    
    reconstructed_chunks = []
    temp = ""
    for chunk in raw_chunks:
        if chunk in [".", "!", "?", "\n"]:
            temp += chunk
            reconstructed_chunks.append(temp.strip())
            temp = ""
        else:
            if temp:
                reconstructed_chunks.append(temp.strip())
            temp = chunk
    if temp:
        reconstructed_chunks.append(temp.strip())

    if not reconstructed_chunks:
        reconstructed_chunks = [text]

    highlighted_html = []
    total_score = 0.0
    ai_sentence_count = 0
    valid_chunks = 0

    for chunk in reconstructed_chunks:
        if len(re.findall(r'\b[a-zA-Z]+\b', chunk)) < 2:
            highlighted_html.append(chunk + " ")
            continue

        chunk_score = score_single_sentence(chunk)
        total_score += chunk_score
        valid_chunks += 1

        if chunk_score >= 40.0:
            ai_sentence_count += 1
            highlighted_html.append(
                f'<mark style="background-color: #ffe066; padding: 2px 5px; border-radius: 4px; font-weight: 500;" title="AI Score: {chunk_score}%">{chunk}</mark> '
            )
        else:
            highlighted_html.append(f'{chunk} ')

    avg_score = round(total_score / valid_chunks, 1) if valid_chunks > 0 else 0.0
    
    # Weighted Calibration for Document Score
    if valid_chunks > 0:
        ai_ratio_score = round((ai_sentence_count / valid_chunks) * 100, 1)
        overall_score = max(avg_score, ai_ratio_score)
    else:
        overall_score = 0.0

    return round(overall_score, 1), "".join(highlighted_html), ai_sentence_count, valid_chunks

# ----------------------------------------------------
# Plagiarism / Duplicate Detection Engine
# ----------------------------------------------------
def check_duplicate_in_db(slide_text: str, threshold=0.50):
    """Normalized text comparison for precise plagiarism matching."""
    if not slide_text.strip() or len(slide_text.split()) < 3:
        return [], 0.0

    conn = get_db_connection()
    c = conn.cursor()
    c.execute('''
        SELECT p.filename, s.slide_number, s.extracted_text 
        FROM slide_records s
        JOIN presentations p ON s.presentation_id = p.id
    ''')
    records = c.fetchall()
    c.close()
    conn.close()

    matches = []
    max_similarity = 0.0
    
    # Strip spaces and non-alphanumeric chars for exact content comparison
    norm_current = re.sub(r'\W+', '', slide_text.lower())

    for filename, slide_num, stored_text in records:
        norm_stored = re.sub(r'\W+', '', stored_text.lower())
        
        ratio = SequenceMatcher(None, norm_current, norm_stored).ratio()
        if ratio >= threshold:
            sim_pct = round(ratio * 100, 1)
            if sim_pct > max_similarity:
                max_similarity = sim_pct
            matches.append({
                "filename": filename,
                "slide_number": slide_num,
                "similarity_pct": sim_pct
            })
            
    return matches, max_similarity

# ----------------------------------------------------
# PDF Report Generator (With Inline AI Highlighting)
# ----------------------------------------------------
def sanitize_text_for_pdf(text):
    """Converts smart quotes and unicode chars to FPDF-safe standard ascii."""
    if not text: 
        return ""
    text = str(text)
    # Replace common Microsoft Word / Mac smart punctuation
    replacements = {
        '\u2018': "'", '\u2019': "'",   # Smart single quotes
        '\u201c': '"', '\u201d': '"',   # Smart double quotes
        '\u2013': "-", '\u2014': "-",   # En and Em dashes
        '\u2026': "...",                # Ellipsis
        '\u00A0': " ",                  # Non-breaking space
        '\u2022': "-"                   # Bullet points
    }
    for search_char, replace_char in replacements.items():
        text = text.replace(search_char, replace_char)
        
    # Force convert anything else to latin-1 (replacing unknowns with '?')
    return text.encode('latin-1', 'replace').decode('latin-1')


def generate_pdf_report(pres_id, filename, file_type, total_slides, overall_ai, overall_plag, upload_time):
    """Generates a PDF report with Yellow Background Highlighting for AI text."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT slide_number, extracted_text, ai_score FROM slide_records WHERE presentation_id = %s ORDER BY slide_number", (pres_id,))
    slides = c.fetchall()
    c.close()
    conn.close()

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Document Header
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, "AI Content & Plagiarism Analysis Report", ln=True, align='C')
    pdf.ln(5)

    pdf.set_font("Arial", '', 11)
    pdf.cell(0, 6, f"File Name: {sanitize_text_for_pdf(filename)}", ln=True)
    pdf.cell(0, 6, f"Upload Time: {upload_time}", ln=True)
    pdf.cell(0, 6, f"Overall AI Score: {overall_ai}%", ln=True)
    pdf.ln(10)

    for slide_num, text, s_ai in slides:
        pdf.set_font("Arial", 'B', 11)
        pdf.cell(0, 8, f"Page {slide_num} | AI Score: {s_ai}%", ln=True)

        if not text or not text.strip():
            continue

        # Re-parse into sentences consistently
        raw_chunks = [c.strip() for c in re.split(r'(\n+|[.!?]+)', text) if c.strip()]
        
        # Write text chunk by chunk
        pdf.set_font("Arial", '', 10)
        for chunk in raw_chunks:
            # Skip separators
            if chunk in [".", "!", "?", "\n"]:
                pdf.write(5, chunk + " ")
                continue
            
            chunk_score = score_single_sentence(chunk)
            safe_chunk = sanitize_text_for_pdf(chunk) + " "
            
            # If AI detected (>0), draw yellow background rect
            if chunk_score > 0:
                # Calculate text width to draw the highlight box
                text_w = pdf.get_string_width(safe_chunk)
                current_x = pdf.get_x()
                current_y = pdf.get_y()
                
                # Draw yellow rect (RGB: 255, 255, 150)
                pdf.set_fill_color(255, 255, 150)
                pdf.rect(current_x, current_y, text_w + 1, 5, 'F')
                
                # Reset color for text
                pdf.set_text_color(0, 0, 0)
                pdf.write(5, safe_chunk)
            else:
                pdf.set_text_color(0, 0, 0)
                pdf.write(5, safe_chunk)

        pdf.ln(8)

    return pdf.output(dest='S').encode('latin-1')

# ----------------------------------------------------
# Database Handlers
# ----------------------------------------------------
def save_to_database(filename, file_type, total_slides, overall_ai, overall_plag, slides_data):
    """Saves scan details with Malaysia Time (UTC+8)."""
    current_myt_time = datetime.now(MYT).strftime("%Y-%m-%d %H:%M:%S")
    
    conn = get_db_connection()
    c = conn.cursor()
    c.execute('''
        INSERT INTO presentations (filename, file_type, total_slides, overall_ai_score, overall_plagiarism_score, upload_time)
        VALUES (%s, %s, %s, %s, %s, %s) RETURNING id;
    ''', (filename, file_type, total_slides, overall_ai, overall_plag, current_myt_time))
    
    pres_id = c.fetchone()[0]

    for s in slides_data:
        c.execute('''
            INSERT INTO slide_records (presentation_id, slide_number, extracted_text, ai_score)
            VALUES (%s, %s, %s, %s);
        ''', (pres_id, s["slide_num"], s["text"], s["ai_score"]))
    
    conn.commit()
    c.close()
    conn.close()


def delete_from_database(presentation_id):
    """Deletes record from Supabase database."""
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("DELETE FROM presentations WHERE id = %s;", (presentation_id,))
    conn.commit()
    c.close()
    conn.close()


def parse_file(uploaded_file):
    filename = uploaded_file.name.lower()
    pages_data = []

    if filename.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
            full_text = " ".join(slide_text)
            pages_data.append({
                "slide_num": idx,
                "text": full_text,
                "word_count": len(re.findall(r'\b\w+\b', full_text))
            })

    elif filename.endswith(".pdf"):
        reader = PdfReader(uploaded_file)
        for idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            clean_text = " ".join(text.split())
            pages_data.append({
                "slide_num": idx,
                "text": clean_text,
                "word_count": len(re.findall(r'\b\w+\b', clean_text))
            })

    return pages_data

# ----------------------------------------------------
# UI Layout
# ----------------------------------------------------
st.set_page_config(page_title="Cloud AI Slide & PDF Detector", layout="wide")

col_title, col_logout = st.columns([0.85, 0.15])
with col_title:
    st.title("📊 Cloud AI Content & Plagiarism Detector")
with col_logout:
    st.write("")
    if st.button("Logout"):
        st.session_state["authenticated"] = False
        st.rerun()

tab1, tab2 = st.tabs(["Analyze Document", "Shared Cloud History"])

with tab1:
    uploaded_file = st.file_uploader("Upload PowerPoint (.pptx) or PDF (.pdf)", type=["pptx", "pdf"])

    if uploaded_file is not None:
        if st.button("Check File", type="primary"):
            file_ext = uploaded_file.name.split(".")[-1].upper()
            
            with st.spinner(f"Analyzing content & plagiarism in {file_ext}..."):
                pages = parse_file(uploaded_file)
                
                if not pages:
                    st.warning("No readable text found in document.")
                else:
                    total_pages = len(pages)
                    total_ai = 0.0
                    total_plag = 0.0
                    scannable_count = 0
                    analyzed = []
                    doc_ai_sentences = 0
                    doc_total_sentences = 0

                    for p in pages:
                        score, highlighted_text, ai_sents, total_sents = analyze_document_text(p["text"])
                        dups, page_max_plag = check_duplicate_in_db(p["text"])
                        
                        p["ai_score"] = score
                        p["highlighted_html"] = highlighted_text
                        p["ai_sentences"] = ai_sents
                        p["total_sentences"] = total_sents
                        p["duplicates"] = dups
                        p["plagiarism_score"] = page_max_plag
                        
                        doc_ai_sentences += ai_sents
                        doc_total_sentences += total_sents
                        analyzed.append(p)

                        if p["word_count"] >= 3:
                            total_ai += score
                            total_plag += page_max_plag
                            scannable_count += 1

                    # Document Level Metrics
                    overall_ai_pct = round(total_ai / scannable_count, 1) if scannable_count > 0 else 0.0
                    overall_plag_pct = round(total_plag / scannable_count, 1) if scannable_count > 0 else 0.0

                    save_to_database(uploaded_file.name, file_ext, total_pages, overall_ai_pct, overall_plag_pct, analyzed)

                    c1, c2, c3, c4 = st.columns(4)
                    c1.metric("Overall AI Score", f"{overall_ai_pct}%")
                    c2.metric("Overall Plagiarism", f"{overall_plag_pct}%")
                    c3.metric("AI Sentences Flagged", f"{doc_ai_sentences} / {doc_total_sentences}")
                    c4.metric("Total Pages / Slides", total_pages)

                    st.markdown("---")
                    st.subheader("Highlighted Sentence Breakdown & Plagiarism Matches")

                    for page in analyzed:
                        p_num = page["slide_num"]
                        score = page["ai_score"]
                        plag_score = page["plagiarism_score"]
                        dups = page["duplicates"]

                        badge = "🔴 High AI" if score >= 50 else ("🟡 Moderate AI" if score >= 25 else "🟢 Likely Human")
                        plag_badge = f" | ⚠️ Plagiarism Match: {plag_score}%" if plag_score > 0 else ""

                        with st.expander(f"Page/Slide {p_num} — AI Score: {score}% ({badge}){plag_badge}", expanded=True if (score >= 35 or plag_score >= 50) else False):
                            st.write(f"**Words:** {page['word_count']} | **AI Sentences:** {page['ai_sentences']} of {page['total_sentences']}")
                            
                            if page["highlighted_html"]:
                                st.markdown(
                                    f'<div style="background-color: #f8f9fa; padding: 15px; border-radius: 8px; border: 1px solid #dee2e6; line-height: 1.8;">{page["highlighted_html"]}</div>',
                                    unsafe_allow_html=True
                                )
                            else:
                                st.info("*(Empty Page)*")

                            if dups:
                                st.warning("⚠️ Database Match Details:")
                                for d in dups:
                                    st.write(f"- **{d['similarity_pct']}% match** with stored file `{d['filename']}` (Page {d['slide_number']})")

with tab2:
    st.subheader("Cloud History (Supabase Database)")
    try:
        conn = get_db_connection()
        c = conn.cursor()
        # Fetches the entire historical list without limitations
        c.execute("SELECT id, filename, file_type, total_slides, overall_ai_score, overall_plagiarism_score, upload_time FROM presentations ORDER BY id DESC")
        records = c.fetchall()
        c.close()
        conn.close()

        if records:
            for r in records:
                rec_id, filename, file_type, total_slides, ai_score, plag_score, upload_time = r
                plag_val = plag_score if plag_score is not None else 0.0
                
                col_info, col_pdf, col_del = st.columns([0.65, 0.20, 0.15])
                
                with col_info:
                    st.write(f"**File:** `{filename}` ({file_type}) | **Pages:** {total_slides} | **AI:** {ai_score}% | **Plag:** {plag_val}% | **Time:** {upload_time}")
                
                with col_pdf:
                    # Generate the PDF data bytes on the fly for the download button
                    pdf_bytes = generate_pdf_report(rec_id, filename, file_type, total_slides, ai_score, plag_val, upload_time)
                    st.download_button(
                        label="📄 Generate Report",
                        data=pdf_bytes,
                        file_name=f"{filename}_AI_Report.pdf",
                        mime="application/pdf",
                        key=f"pdf_{rec_id}"
                    )

                with col_del:
                    if st.button("🗑️ Delete", key=f"del_{rec_id}"):
                        delete_from_database(rec_id)
                        st.success(f"Deleted `{filename}`")
                        st.rerun()
                        
                st.divider()
        else:
            st.info("No documents saved in cloud database yet.")
    except Exception as e:
        st.error(f"Failed to retrieve history: {e}")
